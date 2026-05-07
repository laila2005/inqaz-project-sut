import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()

# Define slide layouts
TITLE_SLIDE_LAYOUT = 0
BULLET_SLIDE_LAYOUT = 1
BLANK_SLIDE_LAYOUT = 6
TITLE_AND_CONTENT_LAYOUT = 5

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    # Styling
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x8B, 0x00, 0x00) # Dark Red
    return slide

def add_bullet_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE_LAYOUT])
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    tf = body.text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
    return slide

def add_image_slide(prs, title_text, img_path1, img_path2=None, label1="", label2=""):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    if img_path2:
        if os.path.exists(img_path1):
            slide.shapes.add_picture(img_path1, Inches(0.5), Inches(2), width=Inches(4))
            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(4), Inches(0.5))
            txbox.text_frame.text = label1
        if os.path.exists(img_path2):
            slide.shapes.add_picture(img_path2, Inches(5), Inches(2), width=Inches(4))
            txbox = slide.shapes.add_textbox(Inches(5), Inches(6), Inches(4), Inches(0.5))
            txbox.text_frame.text = label2
    else:
        if os.path.exists(img_path1):
            slide.shapes.add_picture(img_path1, Inches(2), Inches(2), width=Inches(6))
    return slide

# 1. Title Slide
add_title_slide(prs, "Inqaz AI: Emergency Rescue System", "Computer Vision & Image Processing Final Project\nShorouk Academy (SUT)")

# 2. Overview
add_bullet_slide(prs, "Project Overview", [
    "Objective: Build an end-to-end Computer Vision pipeline to automatically detect car crashes from images.",
    "Goal: Simulate automated emergency alerts to the Ministry of Interior (122) and Ambulance (123).",
    "Requirements Met: Real dataset, full preprocessing, custom CNN, manual transfer learning, strict evaluation, and a working deployment."
])

# 3. Dataset
add_bullet_slide(prs, "Dataset Collection & Quality", [
    "Data Source: Kaggle (no built-in libraries used).",
    "Size: 3,000 real images.",
    "Classes: Crash (Class 0) vs Normal Traffic (Class 1).",
    "Validation: Verified raw image integrity to ensure the model learns true features (crumpled metal, accidents) rather than dataset noise."
])

# 4. Preprocessing
add_bullet_slide(prs, "Preprocessing Pipeline", [
    "Resizing: Uniform 224x224 pixels for CNN and MobileNetV2 compatibility.",
    "Normalization: Scaled pixel values to [-1.0, 1.0] using tf.keras.layers.Rescaling.",
    "Augmentation: Applied Random Flip, 20% Rotation, and 20% Zoom to training data to prevent overfitting.",
    "Splitting: 70% Train (2,100), 15% Val (450), 15% Test (450)."
])

# 5. Scratch CNN
add_bullet_slide(prs, "Model 1: Custom Scratch CNN", [
    "Architecture: 4 Convolutional Blocks (Conv2D -> BatchNorm -> ReLU -> MaxPooling).",
    "Head: Flatten -> Dense(128) + Dropout -> Dense(1) Sigmoid.",
    "Parameters: ~2.5 Million.",
    "Performance: Struggled with the complexity of real-world dashcam images, achieving 63% accuracy."
])

# 6. Transfer Learning
add_bullet_slide(prs, "Model 2: Transfer Learning (MobileNetV2)", [
    "Backbone: MobileNetV2 (ImageNet weights, include_top=False).",
    "Architecture Upgrade: Replaced Flatten() with GlobalAveragePooling2D().",
    "Impact: Reduced parameters from ~16 million to ~1,280, effectively preventing overfitting.",
    "Custom Head: GlobalAveragePooling -> Dense(256) -> Dense(128) -> Sigmoid.",
    "Training: 2-Phase Fine-Tuning (Frozen base -> Unfroze top 20 layers)."
])

# 7. Model Comparison
add_bullet_slide(prs, "Evaluation & Comparison", [
    "Scratch CNN Test Accuracy: 63%",
    "Transfer Learning Test Accuracy: 68%",
    "Crash F1-Score: 0.62 (Scratch) vs 0.68 (Transfer Learning)",
    "Conclusion: MobileNetV2 significantly outperformed the custom model due to its pre-trained feature extraction capabilities and optimized pooling layer."
])

# 8. Confusion Matrices
add_image_slide(prs, "Confusion Matrices", 
                "results/Scratch_CNN_confusion_matrix.png", 
                "results/Transfer_Learning_confusion_matrix.png",
                "Scratch CNN", "Transfer Learning")

# 9. ROC Curves
add_image_slide(prs, "ROC & AUC Curves", 
                "results/Scratch_CNN_roc_curve.png", 
                "results/Transfer_Learning_roc_curve.png",
                "Scratch CNN", "Transfer Learning")

# 10. Deployment
add_bullet_slide(prs, "Deployment (Bonus Phase)", [
    "Platform: Deployed as a web application using Streamlit.",
    "Functionality: Accepts image uploads, applies exact [-1, 1] preprocessing, and runs inference.",
    "Output: Displays crash confidence percentage and simulates emergency dispatches.",
    "Status: Fully operational on localhost."
])

# 11. Conclusion
add_title_slide(prs, "Thank You", "Questions & Discussion")

# Save Presentation
output_path = "Inqaz_Final_Presentation.pptx"
prs.save(output_path)
print(f"Successfully generated {output_path}")
